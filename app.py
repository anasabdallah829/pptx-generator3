import streamlit as st
import zipfile
import os
import io
import shutil
import json
import hashlib
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER, MSO_SHAPE_TYPE
from pptx.util import Inches
import random
from datetime import datetime, date
from PIL import Image
from PIL.ExifTags import TAGS
import tempfile
import streamlit.components.v1 as components

class ConfigManager:
    def __init__(self):
        self.config_file = "template_settings.json"
        
    def save_config(self, pptx_data, config):
        try:
            template_hash = hashlib.md5(pptx_data).hexdigest()
            all_configs = self.load_all_configs()
            all_configs[template_hash] = {
                'timestamp': datetime.now().isoformat(),
                'config': config
            }
            with open(self.config_file, 'w', encoding='utf-8') as f:
                json.dump(all_configs, f, ensure_ascii=False)
            return True
        except Exception:
            return False
            
    def load_config(self, pptx_data):
        try:
            template_hash = hashlib.md5(pptx_data).hexdigest()
            all_configs = self.load_all_configs()
            if template_hash in all_configs:
                return all_configs[template_hash]['config']
        except Exception:
            return None
        return None

    def has_saved_config(self, pptx_data):
        try:
            template_hash = hashlib.md5(pptx_data).hexdigest()
            all_configs = self.load_all_configs()
            return template_hash in all_configs
        except Exception:
            return False

    def load_all_configs(self):
        try:
            if os.path.exists(self.config_file):
                with open(self.config_file, 'r', encoding='utf-8') as f:
                    return json.load(f)
        except Exception:
            pass
        return {}

def init_session():
    defaults = {
        'current_step': 1,
        'pptx_data': None,
        'slide_analysis': None,
        'placeholders_config': {},
        'processing_details': [],
        'show_details_needed': False,
        'selected_placeholder': None,
    }
    for key, value in defaults.items():
        if key not in st.session_state:
            st.session_state[key] = value

def add_detail(message, detail_type="info"):
    st.session_state.processing_details.append({'message': message, 'type': detail_type})
    if detail_type in ['error', 'warning']:
        st.session_state.show_details_needed = True

def clear_details():
    st.session_state.processing_details = []
    st.session_state.show_details_needed = False

def show_details_section():
    if st.session_state.processing_details:
        with st.expander("📋 تفاصيل المعالجة", expanded=st.session_state.show_details_needed):
            for detail in st.session_state.processing_details:
                if detail['type'] == 'success':
                    st.success(detail['message'])
                elif detail['type'] == 'warning':
                    st.warning(detail['message'])
                elif detail['type'] == 'error':
                    st.error(detail['message'])
                else:
                    st.info(detail['message'])

def analyze_slide_placeholders(prs):
    if len(prs.slides) == 0:
        return None
    first_slide = prs.slides[0]
    slide_width, slide_height = prs.slide_width, prs.slide_height
    placeholders = {
        'image_placeholders': [],
        'text_placeholders': [],
        'title_placeholders': [],
        'slide_dimensions': {
            'width': slide_width,
            'height': slide_height,
            'width_inches': slide_width / 914400,
            'height_inches': slide_height / 914400
        }
    }
    placeholder_id = 0

    def clamp_percent(val):
        return max(0, min(val, 100))

    for shape in first_slide.shapes:
        if shape.is_placeholder:
            placeholder_type = shape.placeholder_format.type
            left_percent = clamp_percent((shape.left / slide_width) * 100)
            top_percent = clamp_percent((shape.top / slide_height) * 100)
            width_percent = clamp_percent((shape.width / slide_width) * 100)
            height_percent = clamp_percent((shape.height / slide_height) * 100)
            placeholder_info = {
                'id': placeholder_id,
                'type': placeholder_type,
                'left': shape.left,
                'top': shape.top,
                'width': shape.width,
                'height': shape.height,
                'left_percent': left_percent,
                'top_percent': top_percent,
                'width_percent': width_percent,
                'height_percent': height_percent,
                'rotation': getattr(shape, 'rotation', 0)
            }
            
            if placeholder_type == PP_PLACEHOLDER.PICTURE:
                placeholder_info['current_content'] = "صورة"
                placeholders['image_placeholders'].append(placeholder_info)
            elif placeholder_type == PP_PLACEHOLDER.TITLE:
                placeholder_info['current_content'] = (
                    shape.text_frame.text if hasattr(shape, 'text_frame') and shape.text_frame.text else "العنوان"
                )
                placeholders['title_placeholders'].append(placeholder_info)
            else:
                if hasattr(shape, 'text_frame') and shape.text_frame:
                    placeholder_info['current_content'] = (
                        shape.text_frame.text if shape.text_frame.text else f"نص {placeholder_id + 1}"
                    )
                    placeholders['text_placeholders'].append(placeholder_info)
            placeholder_id += 1

    for shape in first_slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and not shape.is_placeholder:
            left_percent = clamp_percent((shape.left / slide_width) * 100)
            top_percent = clamp_percent((shape.top / slide_height) * 100)
            width_percent = clamp_percent((shape.width / slide_width) * 100)
            height_percent = clamp_percent((shape.height / slide_height) * 100)
            
            image_info = {
                'id': placeholder_id,
                'type': 'regular_image',
                'left': shape.left,
                'top': shape.top,
                'width': shape.width,
                'height': shape.height,
                'left_percent': left_percent,
                'top_percent': top_percent,
                'width_percent': width_percent,
                'height_percent': height_percent,
                'rotation': getattr(shape, 'rotation', 0),
                'current_content': "صورة موجودة"
            }
            placeholders['image_placeholders'].append(image_info)
            placeholder_id += 1
    return placeholders

def render_slide_preview_interactive(slide_analysis):
    dimensions = slide_analysis['slide_dimensions']
    max_width = 820
    aspect_ratio = dimensions['width'] / dimensions['height']
    display_width = max_width if aspect_ratio > 1 else max_width * aspect_ratio
    display_height = max_width / aspect_ratio if aspect_ratio > 1 else max_width

    def generate_placeholder_html(placeholders, type_name, color, icon):
        html = ""
        for i, ph in enumerate(placeholders):
            left = (ph['left_percent'] / 100) * display_width
            top = (ph['top_percent'] / 100) * display_height
            width = (ph['width_percent'] / 100) * display_width
            height = (ph['height_percent'] / 100) * display_height
            
            html += f"""
            <button onclick="window.parent.postMessage({{'phType':'{type_name}','phIdx':{i}}}, '*')" 
                style="
                    position: absolute;
                    left: {left}px;
                    top: {top}px;
                    width: {width}px;
                    height: {height}px;
                    border: 2px solid {color};
                    background: {color}22;
                    font-size: 15px;
                    color: {color};
                    font-weight: bold;
                    border-radius: 5px;
                    cursor: pointer;
                    z-index: 4;
                    transition: all 0.3s ease;
                "
                onmouseover="this.style.background='{color}35'"
                onmouseout="this.style.background='{color}22'"
            >
                {icon} {type_name} {i+1}
            </button>
            """
        return html

    placeholder_html = (
        generate_placeholder_html(slide_analysis['image_placeholders'], 'image', '#ff6b6b', '🖼️') +
        generate_placeholder_html(slide_analysis['text_placeholders'], 'text', '#4ecdc4', '📝') +
        generate_placeholder_html(slide_analysis['title_placeholders'], 'title', '#45b7d1', '📋')
    )

    html_code = f"""
    <div style="
        width: {display_width}px;
        height: {display_height}px;
        border: 2px solid #ddd;
        position: relative;
        background: linear-gradient(135deg, #f5f7fa 0%, #c3cfe2 100%);
        margin: 20px auto;
        border-radius: 10px;
        box-shadow: 0 4px 6px rgba(0, 0, 0, 0.1);
        overflow: hidden;
        z-index: 5;">
        <div style="
            position: absolute;
            top: 5px;
            left: 5px;
            background: rgba(0,0,0,0.65);
            color: white;
            padding: 4px 10px;
            border-radius: 5px;
            font-size: 12px;
            z-index:10;">
            أبعاد الشريحة: {dimensions['width_inches']:.1f}" × {dimensions['height_inches']:.1f}"
        </div>
        {placeholder_html}
    </div>
    <script>
    window.addEventListener('message', function(e) {{
        const d = e.data;
        if (d && d.phType && d.phIdx !== undefined) {{
            window.parent.streamlit && window.parent.streamlit.setComponentValue({{
                type: d.phType,
                index: d.phIdx
            }});
        }}
    }});
    </script>
    """
    
    components.html(html_code, height=int(display_height)+60, scrolling=False)

def configure_selected_placeholder(analysis):
    selected = st.session_state.get('selected_placeholder')
    if not selected:
        st.info("اختر موضع Placeholder من المعاينة أعلاه")
        return

    ptype, idx = selected
    if ptype == 'image':
        ph = analysis['image_placeholders'][idx]
        st.markdown(f"#### إعداد صورة {idx+1}")
        use_image = st.checkbox("استبدال هذه الصورة", value=True, key=f"use_image_{ph['id']}")
        image_order = st.number_input(
            "ترتيب الصورة",
            min_value=1,
            max_value=len(analysis['image_placeholders']),
            value=idx+1,
            key=f"image_order_{ph['id']}"
        )
        
        images_conf = st.session_state.placeholders_config.get('images', {})
        images_conf[f"image_{ph['id']}"] = {
            'use': use_image,
            'order': image_order,
            'placeholder_info': ph
        }
        st.session_state.placeholders_config['images'] = images_conf

    elif ptype == 'text':
        ph = analysis['text_placeholders'][idx]
        st.markdown(f"#### إعداد نص {idx+1}")
        fill_option = st.radio(
            "كيف تريد ملء هذا النص؟",
            ("بدون تغيير", "نص ثابت", "تاريخ", "اسم المجلد"),
            key=f"text_fill_option_{ph['id']}"
        )
        
        placeholder_config = {'type': fill_option, 'value': None}
        if fill_option == "نص ثابت":
            custom_text = st.text_input(
                "أدخل النص المطلوب:",
                key=f"custom_text_{ph['id']}",
                placeholder="مثال: اسم المشروع، اسم الشركة، إلخ..."
            )
            placeholder_config['value'] = custom_text
        elif fill_option == "تاريخ":
            date_option = st.radio(
                "اختر نوع التاريخ:",
                ("تاريخ اليوم", "تاريخ مخصص"),
                key=f"date_option_{ph['id']}"
            )
            if date_option == "تاريخ اليوم":
                placeholder_config['value'] = "today"
            else:
                custom_date = st.date_input(
                    "اختر التاريخ:",
                    key=f"custom_date_{ph['id']}",
                    value=date.today()
                )
                placeholder_config['value'] = custom_date.strftime('%Y-%m-%d')
        elif fill_option == "اسم المجلد":
            placeholder_config['value'] = "folder_name"

        texts_conf = st.session_state.placeholders_config.get('texts', {})
        texts_conf[f"text_{ph['id']}"] = placeholder_config
        st.session_state.placeholders_config['texts'] = texts_conf

    elif ptype == 'title':
        ph = analysis['title_placeholders'][idx]
        st.markdown(f"#### إعداد عنوان {idx+1}")
        st.info("سيتم تعيين عنوان الشريحة باسم المجلد تلقائياً.")

def get_image_date(image_path):
    try:
        with Image.open(image_path) as img:
            exifdata = img.getexif()
            for tag_id in exifdata:
                tag = TAGS.get(tag_id, tag_id)
                data = exifdata.get(tag_id)
                if tag in ['DateTime', 'DateTimeOriginal', 'DateTimeDigitized']:
                    try:
                        date_obj = datetime.strptime(str(data), '%Y:%m:%d %H:%M:%S')
                        return date_obj.strftime('%Y-%m-%d')
                    except Exception:
                        continue
        return datetime.fromtimestamp(os.path.getmtime(image_path)).strftime('%Y-%m-%d')
    except Exception:
        return datetime.now().strftime('%Y-%m-%d')

def apply_configured_placeholders(slide, folder_path, folder_name, slide_analysis, placeholders_config):
    imgs = sorted([f for f in os.listdir(folder_path) 
                  if f.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.webp'))])
    
    # Handle images
    image_config = placeholders_config.get('images', {})
    for config_key, config in image_config.items():
        if config['use'] and config['order'] and config['order'] <= len(imgs):
            image_path = os.path.join(folder_path, imgs[config['order'] - 1])
            placeholder_info = config['placeholder_info']
            
            # Find matching shape
            target_shapes = [
                shape for shape in slide.shapes
                if ((shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER.PICTURE) or
                    (shape.shape_type == MSO_SHAPE_TYPE.PICTURE and not shape.is_placeholder))
            ]
            
            for shape in target_shapes:
                shape_left_percent = (shape.left / slide_analysis['slide_dimensions']['width']) * 100
                shape_top_percent = (shape.top / slide_analysis['slide_dimensions']['height']) * 100
                
                if (abs(shape_left_percent - placeholder_info['left_percent']) < 5 and
                    abs(shape_top_percent - placeholder_info['top_percent']) < 5):
                    try:
                        with open(image_path, "rb") as img_file:
                            if shape.is_placeholder:
                                shape.insert_picture(img_file)
                            else:
                                original_left = shape.left
                                original_top = shape.top
                                original_width = shape.width
                                original_height = shape.height
                                shape_element = shape._element
                                shape_element.getparent().remove(shape_element)
                                slide.shapes.add_picture(
                                    image_path,
                                    original_left,
                                    original_top,
                                    original_width,
                                    original_height
                                )
                        add_detail(f"✅ تم استبدال الصورة {config['order']}: {os.path.basename(image_path)}", "success")
                        break
                    except Exception as e:
                        add_detail(f"❌ فشل استبدال الصورة: {str(e)}", "error")
    
    # Handle text placeholders
    text_config = placeholders_config.get('texts', {})
    text_shapes = [
        shape for shape in slide.shapes
        if (shape.is_placeholder and
            shape.placeholder_format.type not in [PP_PLACEHOLDER.PICTURE, PP_PLACEHOLDER.TITLE] and
            hasattr(shape, 'text_frame'))
    ]
    
    text_index = 0
    for config_key, config in text_config.items():
        if text_index < len(text_shapes):
            shape = text_shapes[text_index]
            try:
                if config['type'] == "بدون تغيير":
                    pass  # Keep original text
                elif config['type'] == "نص ثابت" and config['value']:
                    shape.text_frame.text = config['value']
                elif config['type'] == "تاريخ":
                    if config['value'] == "today":
                        shape.text_frame.text = datetime.now().strftime('%Y-%m-%d')
                    else:
                        shape.text_frame.text = config['value']
                elif config['type'] == "اسم المجلد":
                    shape.text_frame.text = folder_name
                add_detail(f"✅ تم تطبيق النص: {config['type']}", "success")
            except Exception as e:
                add_detail(f"⚠ خطأ في تطبيق النص: {str(e)}", "warning")
            text_index += 1
    
    # Handle title
    title_shapes = [
        shape for shape in slide.shapes
        if shape.is_placeholder and shape.placeholder_format.type == PP_PLACEHOLDER.TITLE
    ]
    if title_shapes:
        title_shapes[0].text = folder_name
        add_detail(f"✅ تم تحديث العنوان: {folder_name}", "success")

def step1_upload_pptx():
    st.title("🔄 PowerPoint Image & Placeholder Replacer")
    st.markdown("---")
    st.markdown("### 📂 الخطوة 1: رفع ملف PowerPoint")
    
    config_manager = ConfigManager()
    
    uploaded_pptx = st.file_uploader(
        "اختر ملف PowerPoint (.pptx)",
        type=["pptx"],
        key="pptx_uploader",
        help="ارفع ملف PowerPoint الذي تريد استبدال الصور والنصوص فيه"
    )
    
    if uploaded_pptx:
        pptx_data = uploaded_pptx.read()
        has_saved_config = config_manager.has_saved_config(pptx_data)
        
        cols = st.columns([1, 1])
        with cols[0]:
            if st.button("📊 تحليل القالب وإعداد جديد", type="primary"):
                with st.spinner("🔍 جاري تحليل ملف PowerPoint..."):
                    try:
                        st.session_state.pptx_data = pptx_data
                        prs = Presentation(io.BytesIO(pptx_data))
                        slide_analysis = analyze_slide_placeholders(prs)
                        if slide_analysis:
                            st.session_state.slide_analysis = slide_analysis
                            st.session_state.current_step = 2
                            st.rerun()
                        else:
                            st.error("❌ لا توجد شرائح في الملف أو حدث خطأ في التحليل")
                    except Exception as e:
                        st.error(f"❌ خطأ في تحليل الملف: {e}")
        
        with cols[1]:
            if has_saved_config:
                if st.button("⚡ استخدام الإعدادات السابقة", type="secondary"):
                    saved_config = config_manager.load_config(pptx_data)
                    if saved_config:
                        st.session_state.pptx_data = pptx_data
                        st.session_state.placeholders_config = saved_config
                        st.session_state.current_step = 3
                        st.rerun()
                    else:
                        st.error("❌ خطأ في تحميل الإعدادات السابقة")

def step2_configure_placeholders():
    st.title("✅ ⚙️ إعداد Placeholders")
    st.markdown("---")
    
    col1, col2, col3 = st.columns([1, 2, 1])
    with col1:
        if st.button("⬅️ العودة للخطوة السابقة"):
            st.session_state.current_step = 1
            st.rerun()
    with col3:
        if st.button("➡️ المتابعة للمعالجة", type="primary"):
            st.session_state.current_step = 3
            st.rerun()

    st.markdown("### 👁️ معاينة القالب")
    analysis = st.session_state.slide_analysis
    if analysis:
        render_slide_preview_interactive(analysis)
        st.markdown("---")
        configure_selected_placeholder(analysis)
        st.markdown("---")
        
        if st.checkbox("📋 عرض ملخص الإعدادات", value=False):
            st.markdown("### 📋 ملخص الإعدادات الحالية")
            image_config = st.session_state.placeholders_config.get('images', {})
            text_config = st.session_state.placeholders_config.get('texts', {})
            
            if image_config:
                st.markdown("#### 🖼️ إعدادات الصور:")
                for key, config in image_config.items():
                    if config['use']:
                        st.success(f"✅ صورة {config['order']}: سيتم استبدالها بالصورة رقم {config['order']} من كل مجلد")
                    else:
                        st.info(f"⏭️ صورة: لن يتم استبدالها")
            
            if text_config:
                st.markdown("#### 📝 إعدادات النصوص:")
                for key, config in text_config.items():
                    if config['type'] == "بدون تغيير":
                        st.info(f"⏭️ نص: سيبقى كما هو")
                    elif config['type'] == "نص ثابت":
                        st.success(f"✅ نص ثابت: '{config['value']}'")
                    elif config['type'] == "تاريخ":
                        st.success(f"📅 تاريخ: {config['value']}")
                    elif config['type'] == "اسم المجلد":
                        st.success(f"📁 اسم المجلد: سيتم استخدام اسم كل مجلد")

def step3_process_files():
    st.title("🚀 معالجة الملفات")
    st.markdown("---")
    
    col1, col2 = st.columns([1, 1])
    with col1:
        if st.button("⬅️ العودة لإعداد Placeholders"):
            st.session_state.current_step = 2
            st.rerun()
            
    st.markdown("### 📂 رفع ملف الصور")
    uploaded_zip = st.file_uploader(
        "اختر ملف ZIP يحتوي على مجلدات صور",
        type=["zip"],
        key="zip_uploader",
        help="ارفع ملف مضغوط يحتوي على مجلدات، كل مجلد يحتوي على صور لشريحة واحدة"
    )

    with st.expander("📋 ملخص الإعدادات المحددة", expanded=True):
        if st.session_state.placeholders_config:
            image_config = st.session_state.placeholders_config.get('images', {})
            text_config = st.session_state.placeholders_config.get('texts', {})
            
            col1, col2 = st.columns(2)
            with col1:
                st.markdown("#### 🖼️ إعدادات الصور:")
                active_images = sum(1 for config in image_config.values() if config['use'])
                for key, config in image_config.items():
                    if config['use']:
                        st.success(f"✅ صورة {config['order']}: سيتم استبدالها")
                    else:
                        st.info(f"⏭️ صورة: لن يتم استبدالها")
                if active_images == 0:
                    st.warning("⚠️ لم يتم تحديد أي صور للاستبدال")
                    
            with col2:
                st.markdown("#### 📝 إعدادات النصوص:")
                active_texts = sum(1 for config in text_config.values() if config['type'] != "بدون تغيير")
                for key, config in text_config.items():
                    if config['type'] != "بدون تغيير":
                        st.success(f"✅ {config['type']}: {config.get('value', 'تلقائي')}")
                    else:
                        st.info(f"⏭️ نص: سيبقى كما هو")
                if active_texts == 0:
                    st.info("ℹ️ لن يتم تغيير أي نصوص")

    st.markdown("### ⚙️ خيارات إضافية")
    col1, col2 = st.columns(2)
    with col1:
        image_order_option = st.radio(
            "ترتيب الصور في المجلدات:",
            ("بالترتيب الأبجدي", "عشوائي"),
            index=0
        )
    with col2:
        skip_empty_folders = st.checkbox(
            "تخطي المجلدات الفارغة",
            value=True
        )

    if uploaded_zip:
        if st.button("🚀 بدء المعالجة", type="primary"):
            clear_details()
            config_manager = ConfigManager()
            
            with tempfile.TemporaryDirectory() as temp_dir:
                try:
                    with st.spinner("📦 جاري استخراج الملفات..."):
                        zip_bytes = io.BytesIO(uploaded_zip.read())
                        with zipfile.ZipFile(zip_bytes, "r") as zip_ref:
                            for member in zip_ref.namelist():
                                if os.path.isabs(member) or ".." in member:
                                    raise Exception("ZIP contains unsafe paths!")
                            zip_ref.extractall(temp_dir)
                    
                    add_detail("📂 تم استخراج الملف المضغوط بنجاح", "success")
                    
                    folder_paths = []
                    for item in os.listdir(temp_dir):
                        item_path = os.path.join(temp_dir, item)
                        if os.path.isdir(item_path):
                            imgs = [f for f in os.listdir(item_path)
                                   if f.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.webp'))]
                            if imgs or not skip_empty_folders:
                                folder_paths.append(item_path)
                                add_detail(f"📁 المجلد '{item}' يحتوي على {len(imgs)} صورة", "info")
                    
                    if not folder_paths:
                        st.error("❌ لا توجد مجلدات تحتوي على صور في الملف المضغوط")
                        add_detail("❌ لا توجد مجلدات تحتوي على صور", "error")
                        show_details_section()
                        return

                    folder_paths.sort()
                    add_detail(f"✅ تم العثور على {len(folder_paths)} مجلد", "success")
                    
                    prs = Presentation(io.BytesIO(st.session_state.pptx_data))
                    first_slide = prs.slides[0]
                    slide_layout = first_slide.slide_layout
                    
                    total_processed = 0
                    created_slides = 0
                    progress_bar = st.progress(0)
                    status_text = st.empty()
                    
                    for folder_idx, folder_path in enumerate(folder_paths):
                        folder_name = os.path.basename(folder_path)
                        status_text.text(f"🔄 معالجة المجلد {folder_idx + 1}/{len(folder_paths)}: {folder_name}")
                        
                        try:
                                                        new_slide = prs.slides.add_slide(slide_layout)
                            created_slides += 1
                            
                            apply_configured_placeholders(
                                new_slide,
                                folder_path,
                                folder_name,
                                st.session_state.slide_analysis,
                                st.session_state.placeholders_config
                            )
                            
                            imgs = [f for f in os.listdir(folder_path)
                                   if f.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.webp'))]
                            total_processed += len(imgs)
                            add_detail(f"✅ تم إنشاء شريحة للمجلد '{folder_name}' مع {len(imgs)} صورة", "success")
                            
                        except Exception as e:
                            add_detail(f"❌ خطأ في معالجة المجلد {folder_name}: {str(e)}", "error")
                        
                        progress_bar.progress((folder_idx + 1) / len(folder_paths))
                    
                    progress_bar.empty()
                    status_text.empty()
                    
                    if created_slides > 0:
                        # Save configuration for future use
                        config_manager.save_config(
                            st.session_state.pptx_data,
                            st.session_state.placeholders_config
                        )
                        
                        st.success("🎉 تم الانتهاء من المعالجة بنجاح!")
                        
                        # Show statistics
                        col1, col2, col3 = st.columns(3)
                        with col1:
                            st.metric("الشرائح المُضافة", created_slides)
                        with col2:
                            st.metric("المجلدات المُعالجة", len(folder_paths))
                        with col3:
                            st.metric("إجمالي الصور", total_processed)
                        
                        # Save output file
                        output_filename = f"PowerPoint_Updated_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
                        output_buffer = io.BytesIO()
                        prs.save(output_buffer)
                        output_buffer.seek(0)
                        
                        st.download_button(
                            label="⬇️ تحميل الملف المُحدث",
                            data=output_buffer.getvalue(),
                            file_name=output_filename,
                            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                            type="primary"
                        )
                        
                        if st.button("🔄 بدء عملية جديدة"):
                            for key in list(st.session_state.keys()):
                                del st.session_state[key]
                            st.rerun()
                        
                    else:
                        st.error("❌ لم يتم إنشاء أي شرائح جديدة")
                    
                    if st.session_state.show_details_needed:
                        show_details_section()
                    elif st.button("📋 إظهار تفاصيل المعالجة"):
                        show_details_section()
                        
                except Exception as e:
                    st.error(f"❌ خطأ أثناء المعالجة: {str(e)}")
                    add_detail(f"❌ خطأ عام في المعالجة: {str(e)}", "error")
                    show_details_section()

def main():
    st.set_page_config(
        page_title="PowerPoint Image Replacer",
        layout="wide",
        page_icon="🔄"
    )
    
    init_session()
    
    # Add custom CSS
    st.markdown("""
    <style>
    .stExpander > div:first-child { background-color: #f0f2f6; }
    .stMetric { 
        background-color: #ffffff;
        padding: 1rem;
        border-radius: 0.5rem;
        border: 1px solid #e1e5e9;
    }
    .slide-preview {
        border: 2px solid #ddd;
        border-radius: 10px;
        box-shadow: 0 4px 6px rgba(0, 0, 0, 0.1);
    }
    .placeholder-btn {
        transition: all 0.3s ease;
    }
    .placeholder-btn:hover {
        transform: scale(1.02);
    }
    </style>
    """, unsafe_allow_html=True)
    
    # Handle different steps
    if st.session_state.current_step == 1:
        step1_upload_pptx()
    elif st.session_state.current_step == 2:
        step2_configure_placeholders()
    elif st.session_state.current_step == 3:
        step3_process_files()
    
    # Show progress bar at bottom
    st.markdown("---")
    progress_labels = ["📂 رفع القالب", "⚙️ إعداد Placeholders", "🚀 المعالجة"]
    cols = st.columns(3)
    for i, (col, label) in enumerate(zip(cols, progress_labels)):
        with col:
            if i + 1 < st.session_state.current_step:
                st.success(f"✅ {label}")
            elif i + 1 == st.session_state.current_step:
                st.info(f"🔄 {label}")
            else:
                st.write(f"⏳ {label}")

if __name__ == '__main__':
    main()
